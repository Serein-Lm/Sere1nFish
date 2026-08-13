"""Shared remote document and attachment resource helpers.

The downloader is intentionally provider-agnostic. Callers supply an aiohttp
session and remain responsible for domain policy; every redirect is still
validated against the public-URL boundary before bytes are read.
"""
from __future__ import annotations

import asyncio
import io
import mimetypes
import os
import re
import shutil
import ssl
import subprocess
import tempfile
import threading
import time
import zipfile
from dataclasses import dataclass
from pathlib import PurePosixPath
from urllib.parse import unquote, urljoin, urlsplit

import aiohttp
from lxml import html as lxml_html

from api.services.url_security import assert_public_http_url


DEFAULT_MAX_ATTACHMENT_TEXT_CHARS = 200_000
ATTACHMENT_EXTENSIONS = {
    ".7z",
    ".csv",
    ".pdf",
    ".doc",
    ".docx",
    ".dps",
    ".et",
    ".xls",
    ".xlsx",
    ".ppt",
    ".pptx",
    ".rar",
    ".rtf",
    ".txt",
    ".wps",
    ".zip",
}
ATTACHMENT_LABEL_MARKERS = (
    "附件",
    "下载",
    "采购文件",
    "招标文件",
    "投标文件",
    "报名表",
    "申请表",
)
_WHITESPACE_RE = re.compile(r"\s+")
_QUOTED_LINK_RE = re.compile(r"['\"]([^'\"<>]+)['\"]")
_HTML_CHARSET_RE = re.compile(
    rb"(?:charset\s*=|encoding\s*=)\s*['\"]?\s*([A-Za-z0-9._-]+)",
    re.I,
)
_LINKISH_SUFFIX_RE = re.compile(
    r"\.(?:s?html?|pdf|docx?|xlsx?|pptx?|wps|et|dps|txt|csv|rtf|zip|rar|7z)"
    r"(?:[?#].*)?$",
    re.I,
)
_PDF_MAX_PAGES = 200
_PDF_OCR_DEADLINE_SECONDS = 900
_PDF_MIN_TEXT_CHARS_PER_PAGE = 24
_DOCX_OCR_MAX_IMAGES = 80
_DOCX_OCR_MAX_IMAGE_BYTES = 32 * 1024 * 1024
_DOCX_OCR_DEADLINE_SECONDS = 900
_DOCX_RASTER_EXTENSIONS = {".bmp", ".jpeg", ".jpg", ".png", ".tif", ".tiff", ".webp"}
_OFFICE_CONVERSION_TIMEOUT_SECONDS = 300
_OFFICE_CONVERSION_SLOTS = threading.BoundedSemaphore(2)
_OPENSSL_LEGACY_SERVER_CONNECT = int(
    getattr(ssl, "OP_LEGACY_SERVER_CONNECT", 0x00000004)
)


@dataclass(slots=True)
class FetchedResource:
    url: str
    data: bytes
    content_type: str
    filename: str


def create_public_fetch_ssl_context() -> ssl.SSLContext:
    """Create the relaxed TLS context used only for public evidence downloads."""
    context = ssl.create_default_context()
    context.check_hostname = False
    context.verify_mode = ssl.CERT_NONE
    # OpenSSL 3 disables legacy server renegotiation by default. Some public
    # institution CDNs still require it before any HTTP response is available.
    context.options |= _OPENSSL_LEGACY_SERVER_CONNECT
    return context


def html_text_and_links(
    content: str | bytes,
    base_url: str,
    *,
    content_type: str = "",
) -> tuple[str, list[dict[str, str]]]:
    if not content:
        return "", []
    decoded = _decode_html_content(content, content_type)
    try:
        root = lxml_html.fromstring(decoded, base_url=base_url)
    except (TypeError, ValueError):
        try:
            root = lxml_html.fragment_fromstring(decoded, create_parent=True)
        except (TypeError, ValueError):
            text = _WHITESPACE_RE.sub(" ", decoded).strip()
            return text, []
    links = extract_html_links(root, base_url)
    for node in root.xpath("//script|//style|//noscript|//template"):
        parent = node.getparent()
        if parent is not None:
            parent.remove(node)
    return _WHITESPACE_RE.sub(" ", root.text_content()).strip(), links


def _decode_html_content(content: str | bytes, content_type: str = "") -> str:
    """Decode public HTML without letting lxml default UTF-8 bytes to Latin-1."""
    if isinstance(content, str):
        return content
    encodings: list[str] = []
    header_match = re.search(r"charset\s*=\s*['\"]?\s*([\w.-]+)", content_type, re.I)
    if header_match:
        encodings.append(header_match.group(1))
    declared_match = _HTML_CHARSET_RE.search(content[:16_384])
    if declared_match:
        encodings.append(declared_match.group(1).decode("ascii", errors="ignore"))
    encodings.extend(("utf-8-sig", "gb18030"))
    for encoding in dict.fromkeys(value for value in encodings if value):
        try:
            decoded = content.decode(encoding)
        except (LookupError, UnicodeDecodeError):
            continue
        return re.sub(
            r"^\ufeff?\s*<\?xml[^>]*\?>",
            "",
            decoded,
            count=1,
            flags=re.I,
        )
    return content.decode("utf-8", errors="replace")


def extract_html_links(root, base_url: str) -> list[dict[str, str]]:
    """Extract navigable and embedded resources through one shared policy."""
    links: list[dict[str, str]] = []
    seen: set[str] = set()
    nodes = root.xpath(
        "self::a[@href or @data-href or @data-url or @onclick] | "
        "self::iframe[@src or @data-src] | self::embed[@src] | "
        "self::object[@data] | "
        ".//a[@href or @data-href or @data-url or @onclick] | "
        ".//iframe[@src or @data-src] | .//embed[@src] | .//object[@data]"
    )
    for node in nodes:
        candidates = [
            str(node.get(attribute) or "").strip()
            for attribute in (
                "href",
                "data-href",
                "data-url",
                "src",
                "data-src",
                "data",
            )
            if str(node.get(attribute) or "").strip()
        ]
        onclick = str(node.get("onclick") or "")
        for quoted in _QUOTED_LINK_RE.findall(onclick):
            value = quoted.strip()
            if value.startswith(("http://", "https://", "/", "./", "../")) or (
                _LINKISH_SUFFIX_RE.search(value)
            ):
                candidates.append(value)
        label = _WHITESPACE_RE.sub(" ", node.text_content()).strip()
        if not label:
            label = str(
                node.get("title")
                or node.get("aria-label")
                or node.get("download")
                or ""
            ).strip()
        for candidate in candidates:
            if not candidate or candidate.lower().startswith(
                ("javascript:", "data:", "mailto:", "tel:", "#")
            ):
                continue
            absolute = urljoin(base_url, candidate)
            try:
                parsed = urlsplit(absolute)
            except ValueError:
                continue
            if parsed.scheme not in {"http", "https"} or not parsed.hostname:
                continue
            if absolute in seen:
                continue
            seen.add(absolute)
            links.append(
                {
                    "url": absolute,
                    "label": label[:300]
                    or safe_remote_filename(PurePosixPath(parsed.path).name),
                }
            )
    return links


def is_attachment_link(link: dict[str, str]) -> bool:
    url = str(link.get("url") or "")
    label = str(link.get("label") or "")
    url_suffix = PurePosixPath(urlsplit(url).path).suffix.lower()
    label_suffix = PurePosixPath(label).suffix.lower()
    return (
        url_suffix in ATTACHMENT_EXTENSIONS
        or label_suffix in ATTACHMENT_EXTENSIONS
        or any(marker in label for marker in ATTACHMENT_LABEL_MARKERS)
    )


def safe_remote_filename(value: str) -> str:
    raw = str(value or "").encode("utf-8", errors="surrogateescape")
    decoded = ""
    for encoding in ("utf-8", "gb18030"):
        try:
            decoded = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if not decoded:
        decoded = raw.decode("utf-8", errors="replace")
    decoded = PurePosixPath(decoded.replace("\\", "/")).name
    decoded = re.sub(r"[\x00-\x1f\x7f]+", "_", decoded).strip(" .")
    return decoded[-180:]


def filename_from_response(
    url: str,
    headers: aiohttp.typedefs.LooseHeaders,
    content_type: str,
) -> str:
    disposition = str(headers.get("Content-Disposition") or "")
    match = re.search(
        r"filename\*?=(?:UTF-8''|\")?([^\";]+)",
        disposition,
        re.I,
    )
    if match:
        name = unquote(match.group(1).strip().strip('"'))
    else:
        name = unquote(PurePosixPath(urlsplit(url).path).name)
    name = safe_remote_filename(name)
    if not name:
        extension = mimetypes.guess_extension(
            content_type.split(";", 1)[0]
        ) or ".bin"
        name = "attachment" + extension
    return name


async def fetch_resource(
    session: aiohttp.ClientSession,
    url: str,
    *,
    max_bytes: int,
    redirects: int = 5,
) -> FetchedResource:
    current = str(url or "").strip()
    for _ in range(max(0, redirects) + 1):
        await assert_public_http_url(current)
        async with session.get(current, allow_redirects=False) as response:
            if response.status in {301, 302, 303, 307, 308}:
                location = response.headers.get("Location")
                if not location:
                    raise RuntimeError(f"HTTP {response.status} 缺少跳转地址")
                current = urljoin(current, location)
                continue
            response.raise_for_status()
            declared_size = int(response.headers.get("Content-Length") or 0)
            if declared_size > max_bytes:
                raise ValueError(
                    f"远程文件超过 {max_bytes // 1024 // 1024} MiB 限制"
                )
            chunks: list[bytes] = []
            size = 0
            async for chunk in response.content.iter_chunked(128 * 1024):
                size += len(chunk)
                if size > max_bytes:
                    raise ValueError(
                        f"远程文件超过 {max_bytes // 1024 // 1024} MiB 限制"
                    )
                chunks.append(chunk)
            content_type = str(
                response.headers.get("Content-Type")
                or "application/octet-stream"
            )
            final_url = str(response.url)
            return FetchedResource(
                url=final_url,
                data=b"".join(chunks),
                content_type=content_type,
                filename=filename_from_response(
                    final_url,
                    response.headers,
                    content_type,
                ),
            )
    raise RuntimeError("远程地址跳转次数过多")


async def fetch_resource_with_retry(
    session: aiohttp.ClientSession,
    url: str,
    *,
    max_bytes: int,
    attempts: int = 3,
) -> FetchedResource:
    """Retry transient transport errors while returning deterministic 4xx."""
    last_error: Exception | None = None
    for attempt in range(max(1, attempts)):
        try:
            return await fetch_resource(session, url, max_bytes=max_bytes)
        except aiohttp.ClientResponseError as exc:
            last_error = exc
            if exc.status not in {408, 425, 429} and exc.status < 500:
                raise
        except (asyncio.TimeoutError, aiohttp.ClientError) as exc:
            last_error = exc
        if attempt + 1 < attempts:
            await asyncio.sleep(0.5 * (2**attempt))
    if last_error is not None:
        raise last_error
    raise RuntimeError("远程资源读取失败")


def _ocr_pdf_pages(
    data: bytes,
    page_numbers: list[int],
) -> tuple[dict[int, str], list[str]]:
    rasterizer = shutil.which("pdftoppm")
    tesseract = shutil.which("tesseract")
    if not rasterizer or not tesseract:
        missing = [
            name
            for name, executable in (
                ("pdftoppm", rasterizer),
                ("tesseract", tesseract),
            )
            if not executable
        ]
        return {}, [f"运行环境缺少 {'/'.join(missing)}，扫描版 PDF 已归档但未 OCR"]

    results: dict[int, str] = {}
    errors: list[str] = []
    deadline = time.monotonic() + _PDF_OCR_DEADLINE_SECONDS
    with tempfile.TemporaryDirectory(prefix="sere1nfish-pdf-") as directory:
        pdf_path = os.path.join(directory, "source.pdf")
        with open(pdf_path, "wb") as handle:
            handle.write(data)
        for page_number in page_numbers:
            remaining = deadline - time.monotonic()
            if remaining <= 0:
                errors.append(
                    f"PDF OCR 超过 {_PDF_OCR_DEADLINE_SECONDS} 秒上限，"
                    f"从第 {page_number} 页起未处理"
                )
                break
            output_prefix = os.path.join(directory, f"page-{page_number}")
            image_path = output_prefix + ".png"
            try:
                rendered = subprocess.run(
                    [
                        rasterizer,
                        "-f",
                        str(page_number),
                        "-l",
                        str(page_number),
                        "-singlefile",
                        "-png",
                        "-r",
                        "180",
                        pdf_path,
                        output_prefix,
                    ],
                    check=False,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=max(1, min(60, int(remaining))),
                )
                if rendered.returncode or not os.path.exists(image_path):
                    detail = _decode_command_output(rendered.stderr).strip()
                    errors.append(
                        f"PDF 第 {page_number} 页渲染失败: "
                        f"{detail or f'退出码 {rendered.returncode}'}"
                    )
                    continue
                remaining = deadline - time.monotonic()
                recognized = subprocess.run(
                    [
                        tesseract,
                        image_path,
                        "stdout",
                        "-l",
                        "chi_sim+eng",
                        "--psm",
                        "6",
                    ],
                    check=False,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=max(1, min(90, int(remaining))),
                )
                text = _WHITESPACE_RE.sub(
                    " ", _decode_command_output(recognized.stdout)
                ).strip()
                if recognized.returncode:
                    detail = _decode_command_output(recognized.stderr).strip()
                    errors.append(
                        f"PDF 第 {page_number} 页 OCR 失败: "
                        f"{detail or f'退出码 {recognized.returncode}'}"
                    )
                elif text:
                    results[page_number] = text
            except (OSError, subprocess.SubprocessError) as exc:
                errors.append(f"PDF 第 {page_number} 页 OCR 失败: {exc}")
            finally:
                try:
                    os.unlink(image_path)
                except OSError:
                    pass
    return results, errors


def _extract_pdf_text(data: bytes, limit: int) -> tuple[str, str]:
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data), strict=False)
        page_count = len(reader.pages)
        selected_count = min(page_count, _PDF_MAX_PAGES)
        page_text: dict[int, str] = {}
        blank_pages: list[int] = []
        errors: list[str] = []
        for index, page in enumerate(reader.pages[:selected_count], start=1):
            extracted = _WHITESPACE_RE.sub(
                " ", (page.extract_text() or "")
            ).strip()
            if len(extracted) >= _PDF_MIN_TEXT_CHARS_PER_PAGE:
                page_text[index] = extracted
            else:
                if extracted:
                    page_text[index] = extracted
                blank_pages.append(index)
        if blank_pages:
            ocr_text, ocr_errors = _ocr_pdf_pages(data, blank_pages)
            for page_number, recognized in ocr_text.items():
                original = page_text.get(page_number, "")
                page_text[page_number] = _WHITESPACE_RE.sub(
                    " ",
                    "\n".join(
                        value for value in (original, recognized) if value
                    ),
                ).strip()
            errors.extend(ocr_errors)
        if page_count > _PDF_MAX_PAGES:
            errors.append(
                f"PDF 共 {page_count} 页，超过 {_PDF_MAX_PAGES} 页安全上限"
            )
        text = _WHITESPACE_RE.sub(
            " ",
            "\n".join(
                page_text.get(index, "")
                for index in range(1, selected_count + 1)
            ),
        ).strip()[:limit]
        if not text and not errors:
            errors.append("PDF OCR 未识别到可读文本")
        return text, "; ".join(errors)[:2_000]
    except Exception as exc:  # noqa: BLE001
        return "", str(exc)


def _extract_docx_text(data: bytes, limit: int) -> tuple[str, str]:
    try:
        from docx import Document

        document = Document(io.BytesIO(data))
        values = [paragraph.text.strip() for paragraph in document.paragraphs]
        size = sum(len(item) for item in values)
        for table in document.tables:
            for row in table.rows:
                line = " | ".join(cell.text.strip() for cell in row.cells)
                values.append(line)
                size += len(line)
                if size >= limit:
                    break
            if size >= limit:
                break
        text = (
            _WHITESPACE_RE.sub(" ", "\n".join(item for item in values if item))
            .strip()[:limit]
        )
        has_embedded_media = False
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as archive:
                has_embedded_media = any(
                    not item.is_dir()
                    and item.filename.startswith("word/media/")
                    for item in archive.infolist()
                )
        except (OSError, zipfile.BadZipFile):
            pass
        if text and (len(text) >= 80 or not has_embedded_media):
            return text, ""
        image_text, image_error = _ocr_docx_images(data, limit)
        if image_text:
            combined = _WHITESPACE_RE.sub(
                " ",
                "\n".join(value for value in (text, image_text) if value),
            ).strip()[:limit]
            return combined, image_error
        converted_text, converted_error = _extract_office_via_pdf(
            data,
            suffix=".docx",
            limit=limit,
        )
        if converted_text:
            combined = _WHITESPACE_RE.sub(
                " ",
                "\n".join(value for value in (text, converted_text) if value),
            ).strip()[:limit]
            return combined, converted_error
        if text:
            return text, "; ".join(
                value for value in (image_error, converted_error) if value
            )[:2_000]
        return "", "; ".join(
            value for value in (image_error, converted_error) if value
        )[:2_000]
    except Exception as exc:  # noqa: BLE001
        return "", str(exc)


def _ocr_docx_images(data: bytes, limit: int) -> tuple[str, str]:
    """OCR image-only DOCX files while keeping conversion bounded."""
    tesseract = shutil.which("tesseract")
    if not tesseract:
        return "", "运行环境缺少 tesseract，扫描版 DOCX 已归档但未 OCR"

    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            media = [
                item
                for item in archive.infolist()
                if not item.is_dir() and item.filename.startswith("word/media/")
            ]
            raster = [
                item
                for item in media
                if PurePosixPath(item.filename).suffix.lower()
                in _DOCX_RASTER_EXTENSIONS
                and 0 < item.file_size <= _DOCX_OCR_MAX_IMAGE_BYTES
            ]
            unsupported = sorted(
                {
                    PurePosixPath(item.filename).suffix.lower().removeprefix(".")
                    or "unknown"
                    for item in media
                    if item not in raster
                }
            )
            selected = raster[:_DOCX_OCR_MAX_IMAGES]
            if not selected:
                suffixes = ", ".join(unsupported)
                detail = f"（内嵌格式: {suffixes}）" if suffixes else ""
                return "", f"DOCX 未提取到可读文本且无可 OCR 位图{detail}"

            texts: list[str] = []
            errors: list[str] = []
            deadline = time.monotonic() + _DOCX_OCR_DEADLINE_SECONDS
            with tempfile.TemporaryDirectory(prefix="sere1nfish-docx-") as directory:
                for index, item in enumerate(selected, start=1):
                    remaining = deadline - time.monotonic()
                    if remaining <= 0:
                        errors.append(
                            f"DOCX OCR 超过 {_DOCX_OCR_DEADLINE_SECONDS} 秒上限，"
                            f"从第 {index} 张图片起未处理"
                        )
                        break
                    suffix = PurePosixPath(item.filename).suffix.lower() or ".img"
                    image_path = os.path.join(directory, f"image-{index}{suffix}")
                    with open(image_path, "wb") as handle:
                        handle.write(archive.read(item))
                    try:
                        recognized = subprocess.run(
                            [
                                tesseract,
                                image_path,
                                "stdout",
                                "-l",
                                "chi_sim+eng",
                                "--psm",
                                "6",
                            ],
                            check=False,
                            stdout=subprocess.PIPE,
                            stderr=subprocess.PIPE,
                            timeout=max(1, min(90, int(remaining))),
                        )
                    except (OSError, subprocess.SubprocessError) as exc:
                        errors.append(f"DOCX 第 {index} 张图片 OCR 失败: {exc}")
                        continue
                    text = _WHITESPACE_RE.sub(
                        " ", _decode_command_output(recognized.stdout)
                    ).strip()
                    if recognized.returncode:
                        detail = _decode_command_output(recognized.stderr).strip()
                        errors.append(
                            f"DOCX 第 {index} 张图片 OCR 失败: "
                            f"{detail or f'退出码 {recognized.returncode}'}"
                        )
                    elif text:
                        texts.append(text)
                    if sum(len(value) for value in texts) >= limit:
                        break
            if len(raster) > _DOCX_OCR_MAX_IMAGES:
                errors.append(
                    f"DOCX 内嵌图片超过 {_DOCX_OCR_MAX_IMAGES} 张安全上限"
                )
            if unsupported:
                errors.append("未处理内嵌格式: " + ", ".join(unsupported))
            joined = "\n".join(texts)[:limit]
            if not joined and not errors:
                errors.append("DOCX 内嵌图片 OCR 未识别到可读文本")
            return joined, "; ".join(errors)[:2_000]
    except (OSError, zipfile.BadZipFile, RuntimeError) as exc:
        return "", str(exc)


def _extract_office_via_pdf(
    data: bytes,
    *,
    suffix: str,
    limit: int,
) -> tuple[str, str]:
    """Render image/vector-only Office files to PDF, then reuse bounded OCR."""
    converter = shutil.which("soffice") or shutil.which("libreoffice")
    if not converter:
        return "", "运行环境缺少 LibreOffice，Office 扫描件已归档但未 OCR"
    safe_suffix = suffix if suffix.startswith(".") else f".{suffix}"
    with _OFFICE_CONVERSION_SLOTS:
        try:
            with tempfile.TemporaryDirectory(
                prefix="sere1nfish-office-"
            ) as directory:
                source_path = os.path.join(directory, f"source{safe_suffix}")
                pdf_path = os.path.join(directory, "source.pdf")
                profile_path = os.path.join(directory, "profile")
                with open(source_path, "wb") as handle:
                    handle.write(data)
                converted = subprocess.run(
                    [
                        converter,
                        "--headless",
                        "--nologo",
                        "--nodefault",
                        "--nofirststartwizard",
                        f"-env:UserInstallation=file://{profile_path}",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        directory,
                        source_path,
                    ],
                    check=False,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=_OFFICE_CONVERSION_TIMEOUT_SECONDS,
                    env={**os.environ, "HOME": directory},
                )
                if converted.returncode or not os.path.exists(pdf_path):
                    detail = _decode_command_output(
                        converted.stderr or converted.stdout
                    ).strip()
                    return "", (
                        "Office 转 PDF 失败: "
                        f"{detail or f'退出码 {converted.returncode}'}"
                    )[:2_000]
                with open(pdf_path, "rb") as handle:
                    pdf_data = handle.read()
                if len(pdf_data) > 200 * 1024 * 1024:
                    return "", "Office 转换后的 PDF 超过 200 MiB 安全上限"
                return _extract_pdf_text(pdf_data, limit)
        except subprocess.TimeoutExpired:
            return "", (
                f"Office 转 PDF 超过 {_OFFICE_CONVERSION_TIMEOUT_SECONDS} 秒上限"
            )
        except OSError as exc:
            return "", f"Office 转 PDF 失败: {exc}"


def _with_office_render_fallback(
    data: bytes,
    *,
    suffix: str,
    limit: int,
    native_text: str,
    native_error: str,
) -> tuple[str, str]:
    """Use LibreOffice rendering when native parsing reports sparse/no content."""
    if native_text and not native_error:
        return native_text, ""
    rendered_text, rendered_error = _extract_office_via_pdf(
        data,
        suffix=suffix,
        limit=limit,
    )
    if rendered_text:
        combined = _WHITESPACE_RE.sub(
            " ",
            "\n".join(value for value in (native_text, rendered_text) if value),
        ).strip()[:limit]
        return combined, rendered_error
    return native_text, "; ".join(
        value for value in (native_error, rendered_error) if value
    )[:2_000]


def _extract_xlsx_text(data: bytes, limit: int) -> tuple[str, str]:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        values: list[str] = []
        size = 0
        has_cell_text = False
        for worksheet in workbook.worksheets[:30]:
            values.append(f"工作表: {worksheet.title}")
            for row_index, row in enumerate(
                worksheet.iter_rows(values_only=True),
                start=1,
            ):
                if row_index > 3_000:
                    break
                line = " | ".join(
                    str(cell).strip() for cell in row[:100] if cell is not None
                )
                if line:
                    values.append(line)
                    size += len(line)
                    has_cell_text = True
                if size >= limit:
                    break
            if size >= limit:
                break
        workbook.close()
        text = _WHITESPACE_RE.sub(" ", "\n".join(values)).strip()[:limit]
        return text, "" if has_cell_text else "XLSX 未提取到可读单元格内容"
    except Exception as exc:  # noqa: BLE001
        return "", str(exc)


def _decode_command_output(data: bytes) -> str:
    for encoding in ("utf-8", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_legacy_doc_text(data: bytes, limit: int) -> tuple[str, str]:
    executable = shutil.which("antiword")
    if not executable:
        return "", "运行环境缺少 antiword，旧版 DOC 已归档但未解析"
    path = ""
    try:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as handle:
            handle.write(data)
            path = handle.name
        result = subprocess.run(
            [executable, path],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=30,
        )
        text = _WHITESPACE_RE.sub(
            " ",
            _decode_command_output(result.stdout),
        ).strip()[:limit]
        error = _decode_command_output(result.stderr).strip()
        if result.returncode and not text:
            return "", error or f"antiword 退出码 {result.returncode}"
        return text, "" if text else (error or "旧版 DOC 未提取到可读文本")
    except (OSError, subprocess.SubprocessError) as exc:
        return "", str(exc)
    finally:
        if path:
            try:
                os.unlink(path)
            except OSError:
                pass


def _extract_legacy_xls_text(data: bytes, limit: int) -> tuple[str, str]:
    try:
        import xlrd

        workbook = xlrd.open_workbook(file_contents=data, on_demand=True)
        values: list[str] = []
        size = 0
        has_cell_text = False
        for sheet in workbook.sheets()[:30]:
            values.append(f"工作表: {sheet.name}")
            for row_index in range(min(sheet.nrows, 3_000)):
                line = " | ".join(
                    str(sheet.cell_value(row_index, column)).strip()
                    for column in range(min(sheet.ncols, 100))
                    if str(sheet.cell_value(row_index, column)).strip()
                )
                if line:
                    values.append(line)
                    size += len(line)
                    has_cell_text = True
                if size >= limit:
                    break
            if size >= limit:
                break
        workbook.release_resources()
        text = _WHITESPACE_RE.sub(" ", "\n".join(values)).strip()[:limit]
        return text, "" if has_cell_text else "旧版 XLS 未提取到可读单元格内容"
    except Exception as exc:  # noqa: BLE001
        return "", str(exc)


def _extract_pptx_text(data: bytes, limit: int) -> tuple[str, str]:
    try:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        values: list[str] = []
        size = 0
        has_content = False
        for slide_index, slide in enumerate(presentation.slides, start=1):
            values.append(f"幻灯片 {slide_index}")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = str(shape.text or "").strip()
                    if text:
                        values.append(text)
                        size += len(text)
                        has_content = True
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        line = " | ".join(cell.text.strip() for cell in row.cells)
                        if line:
                            values.append(line)
                            size += len(line)
                            has_content = True
                if size >= limit:
                    break
            if size >= limit:
                break
        text = _WHITESPACE_RE.sub(" ", "\n".join(values)).strip()[:limit]
        return text, "" if has_content else "PPTX 未提取到可读文本"
    except Exception as exc:  # noqa: BLE001
        return "", str(exc)


def _extract_legacy_ppt_text(data: bytes, limit: int) -> tuple[str, str]:
    executable = shutil.which("catppt")
    if not executable:
        return "", "运行环境缺少 catppt，旧版 PPT 已归档但未解析"
    path = ""
    try:
        with tempfile.NamedTemporaryFile(suffix=".ppt", delete=False) as handle:
            handle.write(data)
            path = handle.name
        result = subprocess.run(
            [executable, path],
            check=False,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=60,
        )
        text = _WHITESPACE_RE.sub(
            " ", _decode_command_output(result.stdout)
        ).strip()[:limit]
        error = _decode_command_output(result.stderr).strip()
        if result.returncode and not text:
            return "", error or f"catppt 退出码 {result.returncode}"
        return text, "" if text else (error or "旧版 PPT 未提取到可读文本")
    except (OSError, subprocess.SubprocessError) as exc:
        return "", str(exc)
    finally:
        if path:
            try:
                os.unlink(path)
            except OSError:
                pass


def _extract_plain_text(data: bytes, limit: int) -> tuple[str, str]:
    encodings = ["utf-8-sig", "gb18030"]
    if data.startswith((b"\xff\xfe", b"\xfe\xff")):
        encodings.insert(0, "utf-16")
    for encoding in encodings:
        try:
            text = data.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        text = data.decode("utf-8", errors="replace")
    text = _WHITESPACE_RE.sub(" ", text).strip()[:limit]
    return text, "" if text else "文本附件为空"


def _extract_rtf_text(data: bytes, limit: int) -> tuple[str, str]:
    try:
        from striprtf.striprtf import rtf_to_text

        decoded, decode_error = _extract_plain_text(data, max(limit * 2, limit))
        if decode_error:
            return "", decode_error
        text = _WHITESPACE_RE.sub(" ", rtf_to_text(decoded)).strip()[:limit]
        return text, "" if text else "RTF 未提取到可读文本"
    except Exception as exc:  # noqa: BLE001
        return "", str(exc)


def _extract_zip_text(data: bytes, limit: int) -> tuple[str, str]:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            entries = [item for item in archive.infolist() if not item.is_dir()]
            if len(entries) > 1_000:
                return "", "ZIP 附件文件数量超过 1000 个安全上限"
            if sum(item.file_size for item in entries) > 200 * 1024 * 1024:
                return "", "ZIP 附件解压后超过 200 MiB 安全上限"
            values: list[str] = []
            errors: list[str] = []
            remaining = limit
            for item in entries[:100]:
                suffix = PurePosixPath(item.filename).suffix.lower()
                if suffix not in ATTACHMENT_EXTENSIONS - {".zip"}:
                    continue
                nested = archive.read(item)
                text, error, _format = extract_attachment_text(
                    nested,
                    filename=item.filename,
                    content_type=mimetypes.guess_type(item.filename)[0] or "",
                    max_chars=remaining,
                )
                if text:
                    block = f"【压缩包文件 {safe_remote_filename(item.filename)}】\n{text}"
                    values.append(block)
                    remaining -= len(block)
                if error:
                    errors.append(f"{safe_remote_filename(item.filename)}: {error}")
                if remaining <= 0:
                    break
            joined = "\n\n".join(values)[:limit]
            if not joined and not errors:
                errors.append("ZIP 中没有可解析的文档附件")
            return joined, "; ".join(errors)[:2_000]
    except (OSError, zipfile.BadZipFile, RuntimeError) as exc:
        return "", str(exc)


def _openxml_kind(data: bytes) -> str:
    if not data.startswith(b"PK"):
        return ""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            entries = archive.infolist()
            if len(entries) > 5_000:
                raise ValueError("Office 附件压缩条目过多")
            if sum(item.file_size for item in entries) > 100 * 1024 * 1024:
                raise ValueError("Office 附件解压后超过 100 MiB 安全上限")
            names = {item.filename for item in entries}
    except (OSError, zipfile.BadZipFile):
        return ""
    if "word/document.xml" in names:
        return "docx"
    if "xl/workbook.xml" in names:
        return "xlsx"
    if "ppt/presentation.xml" in names:
        return "pptx"
    return ""


def extract_attachment_text(
    data: bytes,
    *,
    filename: str,
    content_type: str,
    max_chars: int = DEFAULT_MAX_ATTACHMENT_TEXT_CHARS,
) -> tuple[str, str, str]:
    limit = max(1_000, min(int(max_chars), 300_000))
    suffix = PurePosixPath(str(filename or "")).suffix.lower()
    normalized_type = str(content_type or "").lower()
    if data.startswith(b"%PDF") or "pdf" in normalized_type or suffix == ".pdf":
        text, error = _extract_pdf_text(data, limit)
        return text, error, "pdf"
    try:
        openxml_kind = _openxml_kind(data)
    except ValueError as exc:
        return "", str(exc), suffix.removeprefix(".") or "openxml"
    if (
        openxml_kind == "docx"
        or suffix == ".docx"
        or "wordprocessingml" in normalized_type
    ):
        text, error = _extract_docx_text(data, limit)
        return text, error, "docx"
    if (
        openxml_kind == "xlsx"
        or suffix == ".xlsx"
        or "spreadsheetml" in normalized_type
    ):
        text, error = _extract_xlsx_text(data, limit)
        text, error = _with_office_render_fallback(
            data,
            suffix=".xlsx",
            limit=limit,
            native_text=text,
            native_error=error,
        )
        return text, error, "xlsx"
    if (
        openxml_kind == "pptx"
        or suffix == ".pptx"
        or "presentationml" in normalized_type
    ):
        text, error = _extract_pptx_text(data, limit)
        text, error = _with_office_render_fallback(
            data,
            suffix=".pptx",
            limit=limit,
            native_text=text,
            native_error=error,
        )
        return text, error, "pptx"
    if suffix == ".doc" or "msword" in normalized_type:
        text, error = _extract_legacy_doc_text(data, limit)
        text, error = _with_office_render_fallback(
            data,
            suffix=".doc",
            limit=limit,
            native_text=text,
            native_error=error,
        )
        return text, error, "doc"
    if suffix == ".xls" or "ms-excel" in normalized_type:
        text, error = _extract_legacy_xls_text(data, limit)
        text, error = _with_office_render_fallback(
            data,
            suffix=".xls",
            limit=limit,
            native_text=text,
            native_error=error,
        )
        return text, error, "xls"
    if suffix == ".ppt" or "ms-powerpoint" in normalized_type:
        text, error = _extract_legacy_ppt_text(data, limit)
        text, error = _with_office_render_fallback(
            data,
            suffix=".ppt",
            limit=limit,
            native_text=text,
            native_error=error,
        )
        return text, error, "ppt"
    if suffix in {".wps", ".et", ".dps"}:
        text, error = _extract_office_via_pdf(
            data,
            suffix=suffix,
            limit=limit,
        )
        return text, error, suffix.removeprefix(".")
    if suffix in {".txt", ".csv"} or normalized_type.startswith("text/plain"):
        text, error = _extract_plain_text(data, limit)
        return text, error, suffix.removeprefix(".") or "text"
    if suffix == ".rtf" or "application/rtf" in normalized_type:
        text, error = _extract_rtf_text(data, limit)
        return text, error, "rtf"
    if suffix == ".zip" or normalized_type in {
        "application/zip",
        "application/x-zip-compressed",
    }:
        text, error = _extract_zip_text(data, limit)
        return text, error, "zip"
    return (
        "",
        "当前格式仅归档，未提取正文",
        suffix.removeprefix(".") or "binary",
    )
