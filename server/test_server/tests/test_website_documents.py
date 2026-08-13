from __future__ import annotations

import asyncio
import ssl
from io import BytesIO
from types import SimpleNamespace

import aiohttp

from api.dao import website_crawl as website_crawl_dao
from api.dao.website_crawl import counters_from_summary
from api.services.source_documents.resources import (
    create_public_fetch_ssl_context,
    extract_attachment_text,
    html_text_and_links,
)


def test_public_fetch_ssl_context_supports_legacy_public_sites() -> None:
    context = create_public_fetch_ssl_context()

    assert context.check_hostname is False
    assert context.verify_mode == ssl.CERT_NONE
    assert context.options & int(
        getattr(ssl, "OP_LEGACY_SERVER_CONNECT", 0x00000004)
    )


def test_final_crawl_counters_are_derived_from_authoritative_summary() -> None:
    counters = counters_from_summary(
        {
            "total_pages": 1323,
            "processed_pages": 1323,
            "documents_archived": 1193,
            "attachments_archived": 305,
            "contacts_found": 1701,
            "findings_upserted": 1724,
            "truncated": False,
        }
    )

    assert counters["processed_pages"] == 1323
    assert counters["documents_archived"] == 1193
    assert counters["attachments_archived"] == 305
    assert counters["findings_upserted"] == 1724
    assert "truncated" not in counters


def test_durable_crawl_summary_counts_persisted_findings() -> None:
    class Cursor:
        async def to_list(self, _length):
            return [
                {
                    "_id": {"status": "archived", "kind": "document"},
                    "count": 2,
                    "attachments": 3,
                    "contacts": 4,
                    "findings": 4,
                }
            ]

    class Collection:
        async def find_one(self, *_args, **_kwargs):
            return {"status": "completed", "summary": {"truncated": False}}

        def aggregate(self, _pipeline):
            return Cursor()

    class Db:
        def __getitem__(self, _name):
            return Collection()

    summary = asyncio.run(
        website_crawl_dao.summarize_task(
            Db(),  # type: ignore[arg-type]
            crawl_task_id="crawl-1",
        )
    )

    assert summary["findings_upserted"] == 4
    assert summary["contacts_found"] == 4
from api.services.source_documents.web import (
    OfficialWebDocumentProvider,
    parse_official_html,
)
from api.services.website_documents import (
    archive_page_status,
    classify_discovered_link,
    resolve_website_collection_policy,
    select_official_seed_urls,
    url_matches_required_path_segments,
)


def test_official_html_extracts_article_and_attachment() -> None:
    parsed = parse_official_html(
        b"""
        <html><head><title>site</title></head><body>
          <div class="left fl">
            <h5>Contact notice</h5><div class="articleCon content">
              <p>Call 010-12345678 before 2026-08-20.</p>
            </div>
            <div class="wzFooter"><a href="files/list.xlsx">Attachment</a></div>
          </div>
        </body></html>
        """,
        url="https://example.gov.cn/notices/1.html",
    )

    assert parsed["title"] == "Contact notice"
    assert "010-12345678" in parsed["text"]
    assert parsed["attachment_links"] == [
        {
            "url": "https://example.gov.cn/notices/files/list.xlsx",
            "label": "Attachment",
        }
    ]


def test_official_html_extracts_embedded_and_data_attachments() -> None:
    parsed = parse_official_html(
        b"""
        <html><body><article>
          <h1>Notice</h1><p>Evidence body with enough stable text for selection.</p>
          <iframe title="PDF evidence" src="files/evidence.pdf"></iframe>
          <a data-href="files/list.csv">Download list</a>
        </article></body></html>
        """,
        url="https://example.gov.cn/notices/1.html",
    )

    assert [item["url"] for item in parsed["attachment_links"]] == [
        "https://example.gov.cn/notices/files/evidence.pdf",
        "https://example.gov.cn/notices/files/list.csv",
    ]


def test_official_html_keeps_attachment_sibling_outside_selected_article() -> None:
    parsed = parse_official_html(
        b"""
        <html><body>
          <article><h1>Notice</h1><p>Stable article body content.</p></article>
          <div class="attachment"><a href="files/contact.docx">Attachment</a></div>
        </body></html>
        """,
        url="https://example.gov.cn/notices/1.html",
    )

    assert parsed["attachment_links"] == [
        {
            "url": "https://example.gov.cn/notices/files/contact.docx",
            "label": "Attachment",
        }
    ]


def test_official_html_ignores_local_editor_resource_paths() -> None:
    parsed = parse_official_html(
        b"""
        <html><body><article>
          <h1>Notice</h1><p>Stable article body content.</p>
          <img src="file:///C:/Users/editor/clip_image.jpg">
          <a href="file:///C:/Users/editor/contact.pdf">Attachment</a>
          <img src="images/evidence.jpg">
        </article></body></html>
        """,
        url="https://example.gov.cn/notices/1.html",
    )

    assert parsed["image_urls"] == [
        "https://example.gov.cn/notices/images/evidence.jpg"
    ]
    assert parsed["attachment_links"] == []


def test_listing_link_extraction_reads_onclick_navigation() -> None:
    _text, links = html_text_and_links(
        b"""
        <html><body>
          <a onclick="window.open('index_2.shtml')">Next</a>
          <a data-url="202608/t20260812_1.html">Notice</a>
        </body></html>
        """,
        "https://example.gov.cn/notices/",
    )

    assert [item["url"] for item in links] == [
        "https://example.gov.cn/notices/index_2.shtml",
        "https://example.gov.cn/notices/202608/t20260812_1.html",
    ]


def test_listing_link_extraction_decodes_chinese_without_meta_charset() -> None:
    text, links = html_text_and_links(
        '<html><a href="contact.html">采购公告与联系方式</a></html>'.encode(),
        "https://example.gov.cn/notices/",
    )

    assert text == "采购公告与联系方式"
    assert links[0]["label"] == "采购公告与联系方式"


def test_listing_link_extraction_honors_response_charset() -> None:
    text, links = html_text_and_links(
        '<html><a href="contact.html">采购公告</a></html>'.encode("gb18030"),
        "https://example.gov.cn/notices/",
        content_type="text/html; charset=gb2312",
    )

    assert text == "采购公告"
    assert links[0]["label"] == "采购公告"


def test_crawl_link_classifier_ignores_navigation_but_keeps_notice_scope() -> None:
    navigation = classify_discovered_link(
        url="https://example.gov.cn/about/",
        label="About us",
        parent_url="https://example.gov.cn/notices/",
        parent_relevant=True,
        depth=2,
        max_depth=5,
    )
    pagination = classify_discovered_link(
        url="https://example.gov.cn/notices/index_2.html",
        label="3",
        parent_url="https://example.gov.cn/notices/",
        parent_relevant=True,
        depth=2,
        max_depth=5,
    )
    document = classify_discovered_link(
        url="https://example.gov.cn/notices/202608/t20260812_1.html",
        label="Notice detail",
        parent_url="https://example.gov.cn/notices/",
        parent_relevant=True,
        depth=2,
        max_depth=5,
    )
    cross_scope_document = classify_discovered_link(
        url="https://example.gov.cn/legal/202608/t20260812_2.html",
        label="Legal statement",
        parent_url="https://example.gov.cn/notices/",
        parent_relevant=True,
        depth=2,
        max_depth=5,
    )

    assert navigation is None
    assert pagination and pagination["kind"] == "index"
    assert document and document["kind"] == "document"
    assert cross_scope_document is None


def test_crawl_link_classifier_archives_direct_official_attachment() -> None:
    attachment = classify_discovered_link(
        url="https://example.gov.cn/notices/files/contact-list.xlsx",
        label="联系方式附件",
        parent_url="https://example.gov.cn/notices/index.html",
        parent_relevant=True,
        depth=2,
        max_depth=5,
    )

    assert attachment is not None
    assert attachment["kind"] == "document"
    assert attachment["canonical_url"].endswith("contact-list.xlsx")


def test_crawl_link_classifier_archives_download_endpoint_by_label() -> None:
    attachment = classify_discovered_link(
        url="https://example.gov.cn/file/download?id=42",
        label="附件下载",
        parent_url="https://example.gov.cn/notices/1.html",
        parent_relevant=True,
        depth=2,
        max_depth=5,
    )

    assert attachment is not None
    assert attachment["kind"] == "document"


def test_deep_crawl_enters_unhinted_official_sections() -> None:
    standard = classify_discovered_link(
        url="https://example.gov.cn/organization/",
        label="机构设置",
        parent_url="https://example.gov.cn/",
        parent_relevant=False,
        depth=1,
        max_depth=5,
    )
    deep = classify_discovered_link(
        url="https://example.gov.cn/organization/",
        label="机构设置",
        parent_url="https://example.gov.cn/",
        parent_relevant=False,
        depth=1,
        max_depth=5,
        broad_official_discovery=True,
    )

    assert standard is None
    assert deep and deep["kind"] == "index"


def test_deep_crawl_still_excludes_generic_embedded_systems() -> None:
    candidate = classify_discovered_link(
        url="https://example.gov.cn/kkfileview/index",
        label="在线文件预览",
        parent_url="https://example.gov.cn/",
        parent_relevant=False,
        depth=1,
        max_depth=8,
        broad_official_discovery=True,
    )

    assert candidate is None


def test_crawl_link_classifier_recognizes_query_and_path_pagination() -> None:
    query_page = classify_discovered_link(
        url="https://example.gov.cn/notices/list.shtml?pageIndex=2",
        label="2",
        parent_url="https://example.gov.cn/notices/list.shtml?pageIndex=1",
        parent_relevant=True,
        depth=2,
        max_depth=5,
    )
    path_page = classify_discovered_link(
        url="https://example.gov.cn/notices/page/3/",
        label="下一页",
        parent_url="https://example.gov.cn/notices/page/2/",
        parent_relevant=True,
        depth=3,
        max_depth=5,
    )
    article = classify_discovered_link(
        url="https://example.gov.cn/notices/detail.shtml?id=42",
        label="项目申报通知",
        parent_url="https://example.gov.cn/notices/list.shtml?page=2",
        parent_relevant=True,
        depth=3,
        max_depth=5,
    )

    assert query_page and query_page["kind"] == "index"
    assert path_page and path_page["kind"] == "index"
    assert article and article["kind"] == "document"


def test_crawl_link_classifier_enters_research_and_cooperation_sections() -> None:
    research = classify_discovered_link(
        url="https://example.gov.cn/research/",
        label="科研与学术动态",
        parent_url="https://example.gov.cn/",
        parent_relevant=False,
        depth=1,
        max_depth=5,
    )
    cooperation = classify_discovered_link(
        url="https://example.gov.cn/cooperation/",
        label="交流合作",
        parent_url="https://example.gov.cn/",
        parent_relevant=False,
        depth=1,
        max_depth=5,
    )

    assert research and research["kind"] == "index"
    assert cooperation and cooperation["kind"] == "index"


def test_contact_attribution_failure_keeps_page_retryable() -> None:
    assert archive_page_status(
        {
            "archive_status": "complete",
            "contact_attribution_error": "审核批次超时",
        }
    ) == "partial"
    assert archive_page_status({"archive_status": "complete"}) == "archived"


def test_website_collection_retries_partial_pages_once(monkeypatch) -> None:
    from api.services.website_documents import WebsiteDocumentCollectionService

    service = WebsiteDocumentCollectionService(object())
    results = iter(
        [
            {"status": "partial", "truncated": False},
            {"status": "completed", "truncated": False},
        ]
    )

    async def _run(**_kwargs):
        return next(results)

    monkeypatch.setattr(service, "run", _run)
    result = asyncio.run(service.run_until_stable(project_id="project-1"))

    assert result["status"] == "completed"
    assert result["retry_passes"] == 1


def test_website_discovery_retries_urls_after_enqueue_failure(monkeypatch) -> None:
    from api.services import website_documents
    from api.services.source_documents.resources import FetchedResource
    from api.services.website_documents import (
        WebsiteCollectionPolicy,
        WebsiteDocumentCollectionService,
    )

    seed_url = "https://example.gov.cn/"
    document_url = "https://example.gov.cn/notices/contact.html"
    document_enqueue_attempts = 0
    document_enqueue_batch_sizes: list[int] = []

    async def _noop(*_args, **_kwargs):
        return None

    async def _list_pages(*_args, **_kwargs):
        return []

    async def _enqueue_pages(*_args, pages, **_kwargs):
        nonlocal document_enqueue_attempts
        page = pages[0]
        if page["canonical_url"] == document_url:
            document_enqueue_attempts += 1
            document_enqueue_batch_sizes.append(len(pages))
            if document_enqueue_attempts == 1:
                raise RuntimeError("temporary mongo write failure")
        return [
            website_documents.crawl_dao.page_id_for_url(
                "parent_webdocs",
                item["canonical_url"],
            )
            for item in pages
        ]

    async def _fetch(*_args, **_kwargs):
        return FetchedResource(
            url=seed_url,
            data=(
                '<html><a href="/notices/contact.html">采购公告</a>'
                '<a href="/notices/contact.html">联系方式</a></html>'
            ).encode(),
            content_type="text/html",
            filename="index.html",
        )

    async def _ingest(*_args, **_kwargs):
        return {
            "archive_status": "complete",
            "document_id": "doc-1",
            "version_id": "version-1",
            "attachment_count": 0,
            "contacts": [],
            "fields": {},
        }

    async def _summarize(*_args, **_kwargs):
        return {
            "total_pages": 2,
            "pending_pages": 0,
            "failed_pages": 0,
            "archived_documents": 1,
            "partial_documents": 0,
            "rejected_documents": 0,
            "attachments_archived": 0,
            "contacts_found": 0,
            "findings_upserted": 0,
            "by_status": {"discovered": 1, "archived": 1},
            "by_kind": {"index": 1, "document": 1},
            "truncated": False,
        }

    monkeypatch.setattr(website_documents.crawl_dao, "begin_task", _noop)
    monkeypatch.setattr(website_documents.crawl_dao, "list_pages", _list_pages)
    monkeypatch.setattr(
        website_documents.crawl_dao,
        "enqueue_pages",
        _enqueue_pages,
    )
    monkeypatch.setattr(
        website_documents.crawl_dao,
        "mark_page_started",
        _noop,
    )
    monkeypatch.setattr(
        website_documents.crawl_dao,
        "mark_page_retry",
        _noop,
    )
    monkeypatch.setattr(
        website_documents.crawl_dao,
        "mark_page_terminal",
        _noop,
    )
    monkeypatch.setattr(
        website_documents.crawl_dao,
        "heartbeat_task",
        _noop,
    )
    monkeypatch.setattr(
        website_documents.crawl_dao,
        "summarize_task",
        _summarize,
    )
    monkeypatch.setattr(website_documents.crawl_dao, "finish_task", _noop)
    monkeypatch.setattr(website_documents, "fetch_resource_with_retry", _fetch)
    monkeypatch.setattr(website_documents, "ingest_source_url", _ingest)
    monkeypatch.setattr(website_documents, "update_source_progress", _noop)
    monkeypatch.setattr(
        website_documents.findings_dao,
        "upsert_contact_findings_batch",
        _noop,
    )
    monkeypatch.setattr(
        website_documents.findings_dao,
        "reconcile_contact_findings_for_record",
        _noop,
    )

    from api.services import company_scan_recovery

    monkeypatch.setattr(
        company_scan_recovery,
        "reconcile_terminal_company_website_result",
        _noop,
    )
    service = WebsiteDocumentCollectionService(
        object(),
        policy=WebsiteCollectionPolicy(
            max_pages=20,
            max_documents=10,
            max_depth=3,
            discovery_concurrency=1,
            archive_concurrency=1,
            max_page_attempts=2,
        ),
    )
    result = asyncio.run(
        service.run(
            parent_task_id="parent",
            project_id="project-1",
            target={
                "target_id": "target-1",
                "canonical_name": "测试单位",
                "root_domains": ["example.gov.cn"],
            },
            seed_urls=[seed_url],
        )
    )

    assert document_enqueue_attempts == 2
    assert document_enqueue_batch_sizes == [1, 1]
    assert result["status"] == "completed"
    assert result["documents_archived"] == 1


def test_official_seed_selection_prefers_probed_www_origin() -> None:
    seeds = select_official_seed_urls(
        fallback_urls=["https://example.gov.cn/"],
        known_alive_urls=[
            "http://portal.example.gov.cn:9999/",
            "https://www.example.gov.cn/",
            "https://mail.example.gov.cn/",
        ],
        root_domains=["example.gov.cn"],
    )

    assert seeds == ["https://www.example.gov.cn/"]


def test_official_seed_selection_falls_back_to_apex() -> None:
    seeds = select_official_seed_urls(
        fallback_urls=["example.gov.cn"],
        known_alive_urls=[],
        root_domains=["example.gov.cn"],
    )

    assert seeds == ["https://example.gov.cn/"]


def test_official_seed_selection_isolates_shared_site_path() -> None:
    seeds = select_official_seed_urls(
        fallback_urls=[
            "https://www.10086.cn/",
            "https://www.10086.cn/aboutus/news/pannounce/ah/",
            "https://www.10086.cn/zzxx/ah/",
            "https://www.10086.cn/aboutus/news/pannounce/ha/",
        ],
        known_alive_urls=[
            "https://www.10086.cn/",
            "https://www.10086.cn/aboutus/culture/intro/ah/",
        ],
        root_domains=["10086.cn"],
        required_path_segments=["ah"],
    )

    assert seeds == [
        "https://www.10086.cn/aboutus/culture/intro/ah/",
        "https://www.10086.cn/aboutus/news/pannounce/ah/",
        "https://www.10086.cn/zzxx/ah/",
    ]


def test_required_path_segment_matching_is_exact() -> None:
    assert url_matches_required_path_segments(
        "https://www.10086.cn/aboutus/news/pannounce/ah/index.html",
        ["ah"],
    )
    assert not url_matches_required_path_segments(
        "https://www.10086.cn/aboutus/news/pannounce/ha/index.html",
        ["ah"],
    )


def test_crawl_link_classifier_recognizes_detail_page_patterns() -> None:
    announcement = classify_discovered_link(
        url=(
            "https://www.10086.cn/aboutus/news/pannounce/ah/"
            "index_551_551_detail_55624.html"
        ),
        label="业务升级公告",
        parent_url="https://www.10086.cn/aboutus/news/pannounce/ah/",
        parent_relevant=True,
        depth=1,
        max_depth=5,
        broad_official_discovery=True,
    )
    contact = classify_discovered_link(
        url=(
            "https://www.10086.cn/aboutus/culture/intro/"
            "province_culture_intro_detail/ah/index.html?id=3"
        ),
        label="联系方式",
        parent_url=(
            "https://www.10086.cn/aboutus/culture/intro/"
            "province_culture_intro/ah/"
        ),
        parent_relevant=True,
        depth=1,
        max_depth=5,
        broad_official_discovery=True,
    )

    assert announcement and announcement["kind"] == "document"
    assert contact and contact["kind"] == "document"


def test_deep_website_collection_policy_expands_bounded_budget() -> None:
    tuning = SimpleNamespace(
        website_crawl_max_pages=1200,
        website_crawl_max_documents=400,
        website_crawl_max_depth=5,
        website_crawl_concurrency=12,
        website_archive_concurrency=4,
    )

    standard = resolve_website_collection_policy(tuning, mode="standard")
    deep = resolve_website_collection_policy(tuning, mode="deep")

    assert (standard.max_pages, standard.max_documents, standard.max_depth) == (
        1200,
        400,
        5,
    )
    assert (deep.max_pages, deep.max_documents, deep.max_depth) == (6000, 3000, 8)
    assert deep.discovery_concurrency == 24
    assert deep.archive_concurrency == 8
    assert standard.broad_official_discovery is False
    assert deep.broad_official_discovery is True


def test_truncated_website_crawl_revisits_directory_pages() -> None:
    from api.dao.website_crawl import resumable_page_statuses

    assert "discovered" not in resumable_page_statuses(
        {"summary": {"truncated": False}}
    )
    assert "discovered" in resumable_page_statuses(
        {"summary": {"truncated": True}}
    )
    assert "discovered" in resumable_page_statuses(
        {
            "summary": {"truncated": False},
            "config": {"discovery_policy_version": 1},
        },
        {"discovery_policy_version": 2},
    )
    assert "discovered" in resumable_page_statuses(
        {
            "summary": {"truncated": False},
            "config": {
                "discovery_policy_version": 2,
                "max_depth": 5,
                "broad_official_discovery": False,
            },
        },
        {
            "discovery_policy_version": 2,
            "max_depth": 8,
            "broad_official_discovery": True,
        },
    )


def test_xlsx_attachment_text_is_extracted() -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["contact", "value"])
    sheet.append(["phone", "010-12345678"])
    buffer = BytesIO()
    workbook.save(buffer)
    workbook.close()

    text, error, text_format = extract_attachment_text(
        buffer.getvalue(),
        filename="contacts.xlsx",
        content_type=(
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ),
    )

    assert error == ""
    assert text_format == "xlsx"
    assert "010-12345678" in text


def test_scanned_pdf_uses_bounded_ocr_fallback(monkeypatch) -> None:
    from PIL import Image

    from api.services.source_documents import resources

    buffer = BytesIO()
    Image.new("RGB", (640, 360), "white").save(buffer, format="PDF")
    calls: list[list[int]] = []

    def _ocr(_data: bytes, pages: list[int]):
        calls.append(pages)
        return {1: "项目联系人 张三 010-12345678"}, []

    monkeypatch.setattr(resources, "_ocr_pdf_pages", _ocr)
    text, error, text_format = resources.extract_attachment_text(
        buffer.getvalue(),
        filename="scan.pdf",
        content_type="application/pdf",
    )

    assert calls == [[1]]
    assert error == ""
    assert text_format == "pdf"
    assert "010-12345678" in text


def test_sparse_pdf_page_keeps_native_caption_and_ocr(monkeypatch) -> None:
    from PIL import Image

    from api.services.source_documents import resources

    buffer = BytesIO()
    Image.new("RGB", (640, 360), "white").save(buffer, format="PDF")

    class Page:
        @staticmethod
        def extract_text():
            return "附件一"

    class Reader:
        def __init__(self, *_args, **_kwargs):
            self.pages = [Page()]

    monkeypatch.setattr("pypdf.PdfReader", Reader)
    monkeypatch.setattr(
        resources,
        "_ocr_pdf_pages",
        lambda _data, pages: ({1: "联系人 010-12345678"}, []),
    )

    text, error = resources._extract_pdf_text(buffer.getvalue(), 10_000)

    assert "附件一" in text
    assert "010-12345678" in text
    assert error == ""


def test_pptx_and_plain_text_attachments_are_extracted() -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "项目联络"
    slide.placeholders[1].text = "邮箱 contact@example.gov.cn"
    buffer = BytesIO()
    presentation.save(buffer)

    pptx_text, pptx_error, pptx_format = extract_attachment_text(
        buffer.getvalue(),
        filename="contact.pptx",
        content_type=(
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
    )
    csv_text, csv_error, csv_format = extract_attachment_text(
        "姓名,电话\n张三,010-12345678".encode(),
        filename="contact.csv",
        content_type="text/csv",
    )

    assert pptx_error == ""
    assert pptx_format == "pptx"
    assert "contact@example.gov.cn" in pptx_text
    assert csv_error == ""
    assert csv_format == "csv"
    assert "010-12345678" in csv_text


def test_sparse_office_attachment_uses_rendered_ocr_fallback(monkeypatch) -> None:
    from api.services.source_documents import resources

    monkeypatch.setattr(
        resources,
        "_extract_pptx_text",
        lambda _data, _limit: ("幻灯片 1", "PPTX 未提取到可读文本"),
    )
    monkeypatch.setattr(
        resources,
        "_extract_office_via_pdf",
        lambda _data, **_kwargs: ("联系人 010-12345678", ""),
    )

    text, error, text_format = resources.extract_attachment_text(
        b"not-a-real-pptx",
        filename="scan.pptx",
        content_type=(
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
    )

    assert error == ""
    assert text_format == "pptx"
    assert "010-12345678" in text


def test_wps_family_attachment_uses_office_rendering(monkeypatch) -> None:
    from api.services.source_documents import resources

    calls: list[str] = []

    def _render(_data, *, suffix, limit):
        calls.append(suffix)
        return "联系邮箱 contact@example.gov.cn"[:limit], ""

    monkeypatch.setattr(resources, "_extract_office_via_pdf", _render)

    text, error, text_format = resources.extract_attachment_text(
        b"wps-document",
        filename="notice.wps",
        content_type="application/octet-stream",
    )

    assert calls == [".wps"]
    assert error == ""
    assert text_format == "wps"
    assert "contact@example.gov.cn" in text


def test_permanent_missing_attachment_is_warning(monkeypatch) -> None:
    from api.services.source_documents import web

    async def _missing(*_args, **_kwargs):
        raise aiohttp.ClientResponseError(
            request_info=SimpleNamespace(
                real_url="https://example.gov.cn/gone.pdf"
            ),
            history=(),
            status=404,
            message="Not Found",
        )

    monkeypatch.setattr(web, "fetch_resource_with_retry", _missing)
    attachments, errors, warnings = asyncio.run(
        OfficialWebDocumentProvider._download_attachments(
            object(),
            [{"url": "https://example.gov.cn/gone.pdf", "label": "附件"}],
        )
    )

    assert attachments == []
    assert errors == []
    assert warnings and "HTTP 404" in warnings[0]


def test_attachment_downloads_run_concurrently_and_keep_source_order(
    monkeypatch,
) -> None:
    from api.services.source_documents import web
    from api.services.source_documents.resources import FetchedResource

    active = 0
    max_active = 0

    async def _download(_session, url, **_kwargs):
        nonlocal active, max_active
        active += 1
        max_active = max(max_active, active)
        await asyncio.sleep(0.01)
        active -= 1
        name = url.rsplit("/", 1)[-1]
        return FetchedResource(
            url=url,
            data=f"body-{name}".encode(),
            content_type="text/plain",
            filename=name,
        )

    monkeypatch.setattr(web, "fetch_resource_with_retry", _download)
    monkeypatch.setattr(
        web,
        "extract_attachment_text",
        lambda data, **_kwargs: (data.decode(), "", "txt"),
    )
    links = [
        {"url": f"https://example.gov.cn/{index}.txt", "label": str(index)}
        for index in range(6)
    ]

    attachments, errors, warnings = asyncio.run(
        OfficialWebDocumentProvider._download_attachments(object(), links)
    )

    assert 1 < max_active <= web._ATTACHMENT_DOWNLOAD_CONCURRENCY
    assert [item.index for item in attachments] == list(range(6))
    assert [item.filename for item in attachments] == [
        f"{index}.txt" for index in range(6)
    ]
    assert errors == []
    assert warnings == []


def test_official_provider_captures_direct_attachment(monkeypatch) -> None:
    from api.services.source_documents import web
    from api.services.source_documents.resources import FetchedResource

    requested_limits: list[int] = []

    async def _download(_session, url, *, max_bytes, **_kwargs):
        requested_limits.append(max_bytes)
        return FetchedResource(
            url=url,
            data=b"direct attachment body",
            content_type="application/pdf",
            filename="contact.pdf",
        )

    monkeypatch.setattr(web, "fetch_resource_with_retry", _download)
    monkeypatch.setattr(
        web,
        "extract_attachment_text",
        lambda *_args, **_kwargs: ("联系人 010-12345678", "", "pdf"),
    )

    captured = asyncio.run(
        OfficialWebDocumentProvider().capture(
            "https://example.gov.cn/notices/contact.pdf"
        )
    )

    assert requested_limits == [64 * 1024 * 1024]
    assert captured.text == "联系人 010-12345678"
    assert captured.metadata["direct_attachment"] is True
    assert captured.attachments[0].filename == "contact.pdf"


def test_official_provider_allows_large_download_endpoint(monkeypatch) -> None:
    from api.services.source_documents import web
    from api.services.source_documents.resources import FetchedResource

    requested_limits: list[int] = []

    async def _download(_session, url, *, max_bytes, **_kwargs):
        requested_limits.append(max_bytes)
        return FetchedResource(
            url=url,
            data=b"%PDF direct attachment body",
            content_type="application/pdf",
            filename="contact.pdf",
        )

    monkeypatch.setattr(web, "fetch_resource_with_retry", _download)
    monkeypatch.setattr(
        web,
        "extract_attachment_text",
        lambda *_args, **_kwargs: ("联系人 010-12345678", "", "pdf"),
    )

    captured = asyncio.run(
        OfficialWebDocumentProvider().capture(
            "https://example.gov.cn/file/download?id=42"
        )
    )

    assert requested_limits == [64 * 1024 * 1024]
    assert captured.attachments[0].filename == "contact.pdf"


def test_official_provider_recognizes_download_endpoint_without_file_suffix(
    monkeypatch,
) -> None:
    from api.services.source_documents import web
    from api.services.source_documents.resources import FetchedResource

    async def _download(_session, url, **_kwargs):
        return FetchedResource(
            url=url,
            data=b"%PDF direct attachment body",
            content_type="application/pdf",
            filename="download",
        )

    monkeypatch.setattr(web, "fetch_resource_with_retry", _download)
    monkeypatch.setattr(
        web,
        "extract_attachment_text",
        lambda *_args, **_kwargs: ("联系人 010-12345678", "", "pdf"),
    )

    captured = asyncio.run(
        OfficialWebDocumentProvider().capture(
            "https://example.gov.cn/file/download?id=42"
        )
    )

    assert captured.metadata["direct_attachment"] is True
    assert captured.attachments[0].filename == "download"


def test_permanent_missing_image_is_warning(monkeypatch) -> None:
    from api.services.source_documents import web

    async def _missing(*_args, **_kwargs):
        raise aiohttp.ClientResponseError(
            request_info=SimpleNamespace(
                real_url="https://example.gov.cn/gone.png"
            ),
            history=(),
            status=410,
            message="Gone",
        )

    monkeypatch.setattr(web, "fetch_resource_with_retry", _missing)
    images, errors, warnings = asyncio.run(
        OfficialWebDocumentProvider._download_images(
            object(), ["https://example.gov.cn/gone.png"]
        )
    )

    assert images == []
    assert errors == []
    assert warnings and "HTTP 410" in warnings[0]


def test_official_image_download_allows_high_resolution_source(monkeypatch) -> None:
    from api.services.source_documents import web
    from api.services.source_documents.resources import FetchedResource

    buffer = BytesIO()
    from PIL import Image

    Image.new("RGB", (320, 180), "white").save(buffer, format="JPEG")
    requested_limits: list[int] = []

    async def _download(_session, url, *, max_bytes, **_kwargs):
        requested_limits.append(max_bytes)
        return FetchedResource(
            url=url,
            data=buffer.getvalue(),
            content_type="image/jpeg",
            filename="evidence.jpg",
        )

    monkeypatch.setattr(web, "fetch_resource_with_retry", _download)
    images, errors, warnings = asyncio.run(
        OfficialWebDocumentProvider._download_images(
            object(), ["https://example.gov.cn/evidence.jpg"]
        )
    )

    assert requested_limits == [64 * 1024 * 1024]
    assert len(images) == 1
    assert errors == []
    assert warnings == []
